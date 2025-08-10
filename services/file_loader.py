import os
import requests
import io
import mimetypes
import zipfile
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import pytesseract
from docx import Document
import fitz
import logging
logger = logging.getLogger(__name__)
def extract_text_from_pdf(content: bytes) -> str:
    with fitz.open(stream=content, filetype="pdf") as doc:
        text = "".join(page.get_text() for page in doc)
    return text
def extract_text_from_docx(content: bytes) -> str:
    document = Document(io.BytesIO(content))
    full_text = []
    for para in document.paragraphs:
        full_text.append(para.text)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                full_text.append(cell.text)
    return "\n".join(full_text)
def extract_text_from_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        if slide.has_notes_slide:
            text.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(text)
def extract_text_from_xlsx(content: bytes) -> str:
    wb = load_workbook(io.BytesIO(content), read_only=True)
    text = []
    for sheet in wb.worksheets:
        headers = [cell.value for cell in sheet[1]]
        for row in sheet.iter_rows(min_row=2):
            row_data = [f"{header} is {cell.value}" for header, cell in zip(headers, row) if cell.value]
            text.append(", ".join(row_data))
    return "\n".join(text)
def extract_text_from_image(content: bytes) -> str:
    img = Image.open(io.BytesIO(content))
    return pytesseract.image_to_string(img)
def extract_text_from_zip(content: bytes) -> str:
    full_text = []
    with zipfile.ZipFile(io.BytesIO(content)) as zf:
        for file_info in zf.infolist():
            if file_info.is_dir() or file_info.filename.startswith('__MACOSX'):
                continue
            inner_file_name = file_info.filename
            inner_file_content = zf.read(inner_file_name)
            logger.info(f"Processing '{inner_file_name}' from ZIP archive...")
            try:
                text = dispatch_extraction(inner_file_content, inner_file_name)
                full_text.append(text)
            except ValueError as e:
                logger.warning(f"Skipping unsupported file inside ZIP ('{inner_file_name}'): {e}")
    return "\n\n--- End of File ---\n\n".join(full_text)
def dispatch_extraction(content: bytes, filename: str) -> str:
    file_extension = os.path.splitext(filename)[1].lower()
    if file_extension == '.pdf':
        return extract_text_from_pdf(content)
    elif file_extension == '.docx':
        return extract_text_from_docx(content)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(content)
    elif file_extension == '.xlsx':
        return extract_text_from_xlsx(content)
    elif file_extension in ('.jpeg', '.jpg', '.png'):
        return extract_text_from_image(content)
    elif file_extension == '.zip':
        return extract_text_from_zip(content)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")
def extract_text_from_file(file_url: str) -> str:
    logger.info(f"Downloading file from {file_url}")
    try:
        response = requests.get(file_url)
        response.raise_for_status()
        url_path = file_url.split('?')[0]
        filename = os.path.basename(url_path)
        return dispatch_extraction(response.content, filename)
    except requests.exceptions.RequestException as e:
        logger.error(f"Failed to download {file_url}: {e}")
        raise